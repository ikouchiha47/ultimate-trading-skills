#!/usr/bin/env python3
"""
concall_reader.py — Extract content from a concall. Priority: Transcript PDF > AI Summary > PPT.

Usage:
    # From screener_reader.py --section documents output, pass the concall URLs:
    python concall_reader.py --symbol JUBLFOOD --date "Feb 2026"
    python concall_reader.py --transcript <pdf_url> --ai-summary <screener_url>
    python concall_reader.py --transcript <pdf_url>

    # Run screener_reader first to get URLs, then pass them:
    python screener_reader.py --symbol JUBLFOOD --section documents --output /tmp/docs.json
    python concall_reader.py --from-docs /tmp/docs.json --date "Feb 2026"
"""

from __future__ import annotations

import argparse
import asyncio
import json
import re
import sys
import tempfile
from pathlib import Path

import requests
from pdf_to_md import convert as pdf_convert

BSE_HEADERS = {
    "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10.15; rv:150.0) Gecko/20100101 Firefox/150.0",
    "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
    "Accept-Language": "en-US,en;q=0.9",
    "Connection": "keep-alive",
    "Upgrade-Insecure-Requests": "1",
    "Sec-Fetch-Dest": "document",
    "Sec-Fetch-Mode": "navigate",
    "Sec-Fetch-Site": "cross-site",
    "Pragma": "no-cache",
    "Cache-Control": "no-cache",
}

SESSION_FILE = Path.home() / ".screener_session.json"


def download_pdf(url: str) -> bytes | None:
    """Unified fetch (requests → auto Playwright fallback). BSE transcript PDFs live on the
    Akamai-gated AttachHis host; if even the browser route can't get the PDF, the CALLER cascades
    to the AI summary / PPT (run() priority chain)."""
    try:
        import sys as _sys
        from pathlib import Path as _P
        _root = _P(__file__).resolve().parents[4]          # repo root
        if str(_root) not in _sys.path:
            _sys.path.insert(0, str(_root))
        from framework.fetch import fetch
        data = fetch(url, expect="pdf", referer="https://www.bseindia.com/")
        if data:
            return data
        print("[concall] transcript PDF inaccessible (Akamai) — will try other formats")
    except Exception as e:  # noqa: BLE001
        print(f"[concall] download error: {e}")
    return None


async def fetch_ai_summary(url: str) -> str:
    """Fetch AI summary HTML fragment from screener (requires session)."""
    from playwright.async_api import async_playwright

    if not SESSION_FILE.exists():
        print("[ai_summary] No session file. Run screener_login.py first.")
        return ""

    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        context = await browser.new_context(storage_state=str(SESSION_FILE))
        page = await context.new_page()
        resp = await page.goto(url, wait_until="domcontentloaded")
        if resp.status != 200:
            print(f"[ai_summary] {url} -> {resp.status}")
            await browser.close()
            return ""
        # Content is in <article> or direct body
        article = page.locator("article")
        if await article.count():
            text = await article.inner_text()
        else:
            text = await page.locator("body").inner_text()
        await browser.close()
        return text.strip()


def parse_ppt(url: str) -> str | None:
    """Concall PPT deck -> text. BSE serves these as a PDF OR a real .pptx; handle both via the
    unified fetch (Akamai fallback). PDF -> pdf_convert; pptx (OOXML zip) -> python-pptx."""
    import sys as _sys
    from pathlib import Path as _P
    _root = _P(__file__).resolve().parents[4]
    if str(_root) not in _sys.path:
        _sys.path.insert(0, str(_root))
    from framework.fetch import fetch
    data = fetch(url, expect="any", referer="https://www.bseindia.com/")
    if not data:
        return None
    if data[:4] == b"%PDF":
        md, _ = pdf_convert(data)
        return md
    if data[:4] == b"PK\x03\x04":                       # OOXML .pptx
        try:
            import io
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            out = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [sh.text.strip() for sh in slide.shapes
                         if getattr(sh, "has_text_frame", False) and sh.text.strip()]
                if texts:
                    out.append(f"## Slide {i}\n" + "\n".join(texts))
            return "\n\n".join(out) or None
        except Exception as e:  # noqa: BLE001
            print(f"[concall] pptx parse failed: {e}")
            return None
    return None


def transcribe_file(path: str, model: str = "base.en") -> str | None:
    """THIN primitive: transcribe ONE audio file with faster-whisper (local, no API cost).
    The AGENT orchestrates download/convert/split (yt-dlp + ffmpeg) per the SKILL's chunked-STT
    instructions and calls this per chunk; this just turns one clip into text. Needs the `stt` extra."""
    try:
        from faster_whisper import WhisperModel
    except Exception:  # noqa: BLE001
        print("[concall] faster-whisper not installed (uv pip install -e '.[stt]')")
        return None
    wm = WhisperModel(model, device="cpu", compute_type="int8")
    segs, _ = wm.transcribe(path, beam_size=1)
    return " ".join(s.text.strip() for s in segs).strip() or None


def pick_concall(docs: dict, date_filter: str) -> dict | None:
    """Find a concall entry matching a date string like 'Feb 2026'."""
    date_filter = date_filter.strip().lower()
    for entry in docs.get("concalls", []):
        if date_filter in entry.get("date", "").lower():
            return entry
    return None


def run(args):
    result = {
        "date": args.date or "",
        "source": None,
        "text": "",
    }

    transcript_url = args.transcript
    ai_summary_url = args.ai_summary
    ppt_url = getattr(args, "ppt", None)
    rec_url = getattr(args, "rec", None)

    # Load from docs JSON if provided
    if args.from_docs:
        raw = json.loads(Path(args.from_docs).read_text())
        # screener_reader saves under result["documents"]
        docs = raw.get("documents", raw)
        if not args.date:
            print("ERROR: --date required with --from-docs (e.g. 'Feb 2026')")
            sys.exit(1)
        entry = pick_concall(docs, args.date)
        if not entry:
            print(f"No concall found for '{args.date}'. Available:")
            for c in docs.get("concalls", []):
                print(f"  {c['date']}")
            sys.exit(1)
        transcript_url = entry.get("transcript")
        ai_summary_url = entry.get("ai_summary")
        ppt_url = entry.get("ppt")
        rec_url = entry.get("recording") or entry.get("rec")
        result["date"] = entry["date"]
        print(f"[concall] {entry['date']} — transcript={bool(transcript_url)} "
              f"ai_summary={bool(ai_summary_url)} ppt={bool(ppt_url)} rec={bool(rec_url)}")

    # ── Priority 1: Transcript PDF ────────────────────────────────
    if transcript_url:
        print(f"[concall] Downloading transcript PDF...")
        pdf_bytes = download_pdf(transcript_url)
        if pdf_bytes:
            result["source"] = "transcript_pdf"
            markdown, pages = pdf_convert(pdf_bytes)
            result["text"] = markdown
            result["pages"] = len(pages)
            print(f"[concall] Transcript: {len(markdown)} chars, {len(pages)} pages")

    # ── Priority 2: AI Summary modal ─────────────────────────────
    if not result["text"] and ai_summary_url:
        print(f"[concall] Falling back to AI summary...")
        result["source"] = "ai_summary"
        result["text"] = asyncio.run(fetch_ai_summary(ai_summary_url))
        print(f"[concall] AI Summary: {len(result['text'])} chars extracted")

    # ── Priority 3: PPT deck (PDF or .pptx) ──────────────────────
    if not result["text"] and ppt_url:
        print("[concall] Falling back to PPT deck...")
        text = parse_ppt(ppt_url)
        if text:
            result["source"] = "ppt"
            result["text"] = text
            print(f"[concall] PPT: {len(text)} chars extracted")

    # ── Priority 4: Recording (audio/video) — link only here ──────
    # If only a recording exists, capture the link. Transcription is AGENT-ORCHESTRATED (see the
    # equity-research SKILL "Concall recording → STT" workflow): the agent uses yt-dlp + ffmpeg to
    # download and split into 3-min/30s-overlap chunks, runs `concall_reader.py --audio <chunk>` on
    # each, and stitches. Not auto-run in this script (heavy, opt-in, and a multi-step agent task).
    if not result["text"] and rec_url:
        result["source"] = "recording_link"
        result["text"] = (f"Audio/video recording only — see SKILL 'Concall recording → STT' "
                          f"(agent: yt-dlp + ffmpeg split + `--audio` per chunk).\nRecording: {rec_url}")
        print(f"[concall] Recording link captured: {rec_url}")

    if not result["text"]:
        print("[concall] No content retrieved.")
        sys.exit(1)

    # Output
    print(f"\n── SOURCE: {result['source']} ({'%.1f' % (len(result['text'])/1000)}k chars) ──\n")
    preview = result["text"][:3000]
    print(preview)
    if len(result["text"]) > 3000:
        print(f"\n... [{len(result['text']) - 3000} more chars]")

    if args.output:
        out = Path(args.output)
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_text(json.dumps(result, indent=2))
        print(f"\n[saved] {out}")

    return result


def main():
    parser = argparse.ArgumentParser(description="Read concall transcript, AI summary, or PPT")
    parser.add_argument("--from-docs", help="Path to JSON output from screener_reader.py --section documents")
    parser.add_argument("--date", help="Concall date to pick, e.g. 'Feb 2026'")
    parser.add_argument("--transcript", help="Direct BSE transcript PDF URL")
    parser.add_argument("--ai-summary", help="Screener AI summary URL (/concalls/summary/<id>/)")
    parser.add_argument("--audio", help="Transcribe ONE audio chunk (agent-split) via faster-whisper")
    parser.add_argument("--model", default="base.en", help="faster-whisper model (default base.en)")
    parser.add_argument("--output", help="Save JSON/text result to this path")
    args = parser.parse_args()

    # thin STT primitive — the agent splits with ffmpeg and calls this per chunk (see SKILL)
    if args.audio:
        text = transcribe_file(args.audio, args.model) or ""
        if args.output:
            Path(args.output).write_text(text)
            print(f"[saved] {args.output} ({len(text)} chars)")
        else:
            print(text)
        return

    if not args.from_docs and not args.transcript and not args.ai_summary:
        parser.error("Provide --from-docs, --transcript/--ai-summary, or --audio <chunk>")

    run(args)


if __name__ == "__main__":
    main()
